"""Extract full text from all project documents for planning."""
import os, sys

BASE = r"c:\Users\12132\Desktop\DigiFacil\DigiFacil"

try:
    from docx import Document
except ImportError:
    os.system(f"{sys.executable} -m pip install python-docx -q")
    from docx import Document

try:
    from openpyxl import load_workbook
except ImportError:
    os.system(f"{sys.executable} -m pip install openpyxl -q")
    from openpyxl import load_workbook

try:
    from pptx import Presentation
except ImportError:
    os.system(f"{sys.executable} -m pip install python-pptx -q")
    from pptx import Presentation

# --- DOCX files ---
docx_files = [
    "Todo_Directo_Pricing_Strategy.docx",
    "Todo_Directo_Operations_SOP.docx",
    "Todo_Directo_SOP_v2.docx",
    "Todo_Directo_2_Paginas_Sales_Cheat_Sheet_ES_v2.docx",
    "Todo_Directo_Intake_Form.docx",
]

for fname in docx_files:
    path = os.path.join(BASE, fname)
    if not os.path.exists(path):
        print(f"\n=== {fname} === NOT FOUND")
        continue
    print(f"\n{'='*80}")
    print(f"FILE: {fname}")
    print(f"{'='*80}")
    doc = Document(path)
    for i, p in enumerate(doc.paragraphs):
        txt = p.text.strip()
        if txt:
            style = p.style.name if p.style else "None"
            print(f"[{style}] {txt}")
    for t_idx, table in enumerate(doc.tables):
        print(f"\n--- TABLE {t_idx+1} ---")
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            print(" | ".join(cells))

# --- XLSX ---
xlsx_path = os.path.join(BASE, "Todo_Directo_Founder_Financial_Dashboard.xlsx")
print(f"\n{'='*80}")
print(f"FILE: Todo_Directo_Founder_Financial_Dashboard.xlsx")
print(f"{'='*80}")
wb = load_workbook(xlsx_path, data_only=True)
print(f"Sheets: {wb.sheetnames}")
for sname in wb.sheetnames:
    ws = wb[sname]
    print(f"\n--- Sheet: {sname} ({ws.max_row} rows x {ws.max_column} cols) ---")
    for row in ws.iter_rows(min_row=1, max_row=min(ws.max_row, 30), values_only=False):
        vals = [str(c.value) if c.value is not None else "" for c in row]
        print(" | ".join(vals))

# --- PPTX ---
pptx_path = os.path.join(BASE, "Todo_Directo_Pitch_Deck_UPDATED.pptx")
print(f"\n{'='*80}")
print(f"FILE: Todo_Directo_Pitch_Deck_UPDATED.pptx")
print(f"{'='*80}")
prs = Presentation(pptx_path)
for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} ---")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    print(f"  {t}")
        if shape.has_table:
            print("  [TABLE]")
            for row in shape.table.rows:
                cells = [c.text.strip() for c in row.cells]
                print(f"  {' | '.join(cells)}")
