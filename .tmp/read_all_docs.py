import docx

docs = {
    'Pricing Strategy': 'Todo_Directo_Pricing_Strategy.docx',
    'SOP v2': 'Todo_Directo_SOP_v2.docx', 
    'Intake Form': 'Todo_Directo_Intake_Form.docx'
}

for title, fname in docs.items():
    print(f"\n{'='*80}")
    print(f"{title}")
    print(f"{'='*80}")
    try:
        doc = docx.Document(fname)
        for para in doc.paragraphs:
            if para.text.strip():
                print(para.text)
        # Also extract tables
        for i, table in enumerate(doc.tables):
            print(f"\n--- Table {i+1} ---")
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                print(' | '.join(cells))
    except Exception as e:
        print(f"Error: {e}")

# Also read pitch deck slides
print(f"\n{'='*80}")
print("PITCH DECK")
print(f"{'='*80}")
try:
    from pptx import Presentation
    prs = Presentation('Todo_Directo_Pitch_Deck_UPDATED.pptx')
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                print(shape.text)
except Exception as e:
    print(f"Error: {e}")
