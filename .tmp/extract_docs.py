"""Extract text summaries from all project documents."""
import os, sys

BASE = r"c:\Users\12132\Desktop\DigiFacil\DigiFacil"

# --- DOCX ---
try:
    from docx import Document
except ImportError:
    os.system(f"{sys.executable} -m pip install python-docx -q")
    from docx import Document

docx_files = [
    "Todo_Directo_Pricing_Strategy.docx",
    "Todo_Directo_Operations_SOP.docx",
    "Todo_Directo_SOP_v2.docx",
    "Todo_Directo_2_Paginas_Sales_Cheat_Sheet_ES_v2.docx",
    "Todo_Directo_Intake_Form.docx",
]

for fname in docx_files:
    path = os.path.join(BASE, fname)
    if not os.path.exists(path):
        print(f"\n=== {fname} === NOT FOUND")
        continue
    print(f"\n{'='*80}")
    print(f"=== {fname} ===")
    print(f"{'='*80}")
    doc = Document(path)
    lines = []
    for p in doc.paragraphs:
        txt = p.text.strip()
        if txt:
            lines.append(txt)
    # Also check tables
    for t_idx, table in enumerate(doc.tables):
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            lines.append(" | ".join(cells))
    print("\n".join(lines[:200]))  # First 200 lines

# --- XLSX ---
try:
    from openpyxl import load_workbook
except ImportError:
    os.system(f"{sys.executable} -m pip install openpyxl -q")
    from openpyxl import load_workbook

xlsx_path = os.path.join(BASE, "Todo_Directo_Founder_Financial_Dashboard.xlsx")
print(f"\n{'='*80}")
print(f"=== Todo_Directo_Founder_Financial_Dashboard.xlsx ===")
print(f"{'='*80}")
wb = load_workbook(xlsx_path, data_only=True)
print(f"Sheet names: {wb.sheetnames}")
for sname in wb.sheetnames:
    ws = wb[sname]
    print(f"\n--- Sheet: {sname} (rows={ws.max_row}, cols={ws.max_column}) ---")
    for row in ws.iter_rows(min_row=1, max_row=min(ws.max_row, 15), values_only=False):
        vals = [str(c.value) if c.value is not None else "" for c in row]
        print(" | ".join(vals))

# --- PPTX ---
try:
    from pptx import Presentation
except ImportError:
    os.system(f"{sys.executable} -m pip install python-pptx -q")
    from pptx import Presentation

pptx_path = os.path.join(BASE, "Todo_Directo_Pitch_Deck_UPDATED.pptx")
print(f"\n{'='*80}")
print(f"=== Todo_Directo_Pitch_Deck_UPDATED.pptx ===")
print(f"{'='*80}")
prs = Presentation(pptx_path)
for i, slide in enumerate(prs.slides):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    print(f"\nSlide {i+1}: {' // '.join(texts[:5])}")
