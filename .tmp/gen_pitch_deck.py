"""Generate Cresca Pitch Deck — 12 slides with 3-Tier Model."""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0x7B, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
GREEN = RGBColor(0x10, 0xB9, 0x81)
GRAY = RGBColor(0x6B, 0x72, 0x80)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf

def add_bullet(tf, text, font_size=16, color=DARK, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.level = 0
    return p

# ===== SLIDE 1: Title =====
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1, DARK)
add_text_box(s1, 2, 1.5, 9, 1.5, 'CRESCA', 60, True, WHITE, PP_ALIGN.CENTER)
add_text_box(s1, 2, 3.2, 9, 1, 'Sistemas de Crecimiento Empresarial', 28, False, ACCENT, PP_ALIGN.CENTER)
add_text_box(s1, 2, 4.5, 9, 0.8, 'Infraestructura Digital  ·  Automatización  ·  Crecimiento Estructurado', 18, False, WHITE, PP_ALIGN.CENTER)
add_text_box(s1, 2, 5.8, 9, 0.5, 'El Salvador, C.A.', 14, False, GRAY, PP_ALIGN.CENTER)

# ===== SLIDE 2: Positioning =====
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s2, WHITE)
add_text_box(s2, 1, 0.5, 11, 1, 'NO SOMOS UNA AGENCIA.', 40, True, DARK, PP_ALIGN.LEFT)
add_text_box(s2, 1, 1.8, 11, 1, 'SOMOS UNA EMPRESA DE SISTEMAS DE CRECIMIENTO.', 32, True, ACCENT, PP_ALIGN.LEFT)
tf2 = add_text_box(s2, 1, 3.5, 5, 3, 'Lo que NO somos:', 20, True, DARK)
add_bullet(tf2, '✗  Agencia de marketing', 18, GRAY)
add_bullet(tf2, '✗  Constructor de sitios web', 18, GRAY)
add_bullet(tf2, '✗  Agencia de chatbots', 18, GRAY)
add_bullet(tf2, '✗  Servicios sueltos', 18, GRAY)

tf2b = add_text_box(s2, 7, 3.5, 5, 3, 'Lo que SÍ somos:', 20, True, DARK)
add_bullet(tf2b, '✓  Sistemas integrados', 18, GREEN)
add_bullet(tf2b, '✓  Infraestructura digital', 18, GREEN)
add_bullet(tf2b, '✓  Automatización empresarial', 18, GREEN)
add_bullet(tf2b, '✓  Crecimiento estructurado', 18, GREEN)

# ===== SLIDE 3: The Problem =====
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s3, DARK)
add_text_box(s3, 1, 0.5, 11, 1, 'EL PROBLEMA', 40, True, WHITE, PP_ALIGN.LEFT)
add_text_box(s3, 1, 1.8, 11, 0.8, 'Los negocios en El Salvador operan sin sistemas.', 24, False, ACCENT)

tf3 = add_text_box(s3, 1, 3, 5, 4, '', 18, False, WHITE)
add_bullet(tf3, '📋  Procesos 100% manuales', 18, WHITE)
add_bullet(tf3, '📱  Dependen de WhatsApp personal', 18, WHITE)
add_bullet(tf3, '🛒  Solo venden por Facebook Marketplace', 18, WHITE)
add_bullet(tf3, '📉  Pierden clientes por falta de seguimiento', 18, WHITE)
add_bullet(tf3, '⏰  El dueño hace TODO personalmente', 18, WHITE)

tf3b = add_text_box(s3, 7, 3, 5, 4, '', 18, False, WHITE)
add_bullet(tf3b, 'Sin presencia web profesional', 18, GRAY)
add_bullet(tf3b, 'Sin automatización ni workflows', 18, GRAY)
add_bullet(tf3b, 'Sin reportes ni datos', 18, GRAY)
add_bullet(tf3b, 'No escalan: el negocio depende', 18, GRAY)
add_bullet(tf3b, '100% del dueño', 18, GRAY)

# ===== SLIDE 4: The Solution =====
s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s4, WHITE)
add_text_box(s4, 1, 0.5, 11, 1, 'LA SOLUCIÓN: SISTEMAS DE CRECIMIENTO', 36, True, DARK, PP_ALIGN.LEFT)
add_text_box(s4, 1, 1.8, 11, 1, 'Cresca construye la infraestructura digital y automatización que tu negocio necesita para crecer — empezando desde $9.99/mes.', 20, False, GRAY)

tf4 = add_text_box(s4, 1, 3.3, 5, 3.5, '🏗️  Infraestructura Digital', 22, True, DARK)
add_bullet(tf4, 'Presencia web profesional', 16, DARK)
add_bullet(tf4, 'Captura de leads automática', 16, DARK)
add_bullet(tf4, 'Módulo de reservas integrado', 16, DARK)
add_bullet(tf4, 'Google Business optimizado', 16, DARK)

tf4b = add_text_box(s4, 7, 3.3, 5, 3.5, '⚙️  Automatización', 22, True, DARK)
add_bullet(tf4b, 'Workflows automatizados', 16, DARK)
add_bullet(tf4b, 'Seguimiento sin intervención', 16, DARK)
add_bullet(tf4b, 'Respuestas automáticas 24/7', 16, DARK)
add_bullet(tf4b, 'Reportería y optimización', 16, DARK)

# ===== SLIDE 5: Three-Tier Packages =====
s5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s5, LIGHT_BG)
add_text_box(s5, 1, 0.2, 11, 1, 'NUESTROS 3 NIVELES', 36, True, DARK, PP_ALIGN.CENTER)

# STARTER
tf5a = add_text_box(s5, 0.3, 1.3, 4, 5.5, 'STARTER', 26, True, ACCENT)
add_bullet(tf5a, '', 6)
add_bullet(tf5a, '$9.99 – $49.99/mes', 20, ACCENT, True)
add_bullet(tf5a, '', 4)
add_bullet(tf5a, 'Entrada al mundo digital', 14, DARK, True)
add_bullet(tf5a, '', 4)
add_bullet(tf5a, '✓ 1 Workflow Lite incluido', 14, DARK)
add_bullet(tf5a, '✓ Plantillas y guías', 14, DARK)
add_bullet(tf5a, '✓ Google Business (según tier)', 14, DARK)
add_bullet(tf5a, '✓ Micro-sitio (tier Pro)', 14, DARK)
add_bullet(tf5a, '', 4)
add_bullet(tf5a, 'Setup: $0 – $199', 13, GRAY)

# GROWTH
tf5b = add_text_box(s5, 4.6, 1.3, 4, 5.5, 'GROWTH', 26, True, GREEN)
add_bullet(tf5b, '', 6)
add_bullet(tf5b, '$50 – $100/mes', 20, GREEN, True)
add_bullet(tf5b, '', 4)
add_bullet(tf5b, 'Infraestructura + automatización', 14, DARK, True)
add_bullet(tf5b, '', 4)
add_bullet(tf5b, '✓ Website 3–5 páginas', 14, DARK)
add_bullet(tf5b, '✓ Captura de leads + CTA', 14, DARK)
add_bullet(tf5b, '✓ 1–2 Full Workflows', 14, DARK)
add_bullet(tf5b, '✓ Módulo de reservas (si aplica)', 14, DARK)
add_bullet(tf5b, '', 4)
add_bullet(tf5b, 'Setup: $199 – $499', 13, GRAY)

# PRO AUTOMATION
tf5c = add_text_box(s5, 8.9, 1.3, 4, 5.5, 'PRO AUTOMATION', 26, True, ORANGE)
add_bullet(tf5c, '', 6)
add_bullet(tf5c, 'Cotización', 20, ORANGE, True)
add_bullet(tf5c, '', 4)
add_bullet(tf5c, 'Automatización avanzada', 14, DARK, True)
add_bullet(tf5c, '', 4)
add_bullet(tf5c, '✓ Auditoría de procesos', 14, DARK)
add_bullet(tf5c, '✓ 4+ workflows personalizados', 14, DARK)
add_bullet(tf5c, '✓ Monitoreo y optimización', 14, DARK)
add_bullet(tf5c, '✓ Reportería y dashboards', 14, DARK)
add_bullet(tf5c, '', 4)
add_bullet(tf5c, 'Recomendado: 12 meses', 13, GRAY)

# ===== SLIDE 6: Empiece Pequeño, Escale Cuando Crezca =====
s6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s6, DARK)
add_text_box(s6, 1, 0.5, 11, 1, 'EMPIECE PEQUEÑO, ESCALE CUANDO CREZCA', 36, True, WHITE, PP_ALIGN.CENTER)
add_text_box(s6, 1, 1.5, 11, 0.8, 'Su negocio crece con usted — su sistema también.', 20, False, ACCENT, PP_ALIGN.CENTER)

tf6a = add_text_box(s6, 0.5, 2.8, 3.5, 3.5, '📱 STARTER', 24, True, ACCENT)
add_bullet(tf6a, '', 6)
add_bullet(tf6a, 'Desde $9.99/mes', 18, WHITE, True)
add_bullet(tf6a, '1 Workflow Lite', 16, WHITE)
add_bullet(tf6a, 'Plantillas y guías', 16, WHITE)

add_text_box(s6, 4.2, 4, 1, 1, '→', 36, True, ACCENT, PP_ALIGN.CENTER)

tf6b = add_text_box(s6, 5, 2.8, 3.5, 3.5, '🏗️ GROWTH', 24, True, GREEN)
add_bullet(tf6b, '', 6)
add_bullet(tf6b, 'Desde $50/mes', 18, WHITE, True)
add_bullet(tf6b, 'Website + Workflows', 16, WHITE)
add_bullet(tf6b, 'Reservas + Leads', 16, WHITE)

add_text_box(s6, 8.7, 4, 1, 1, '→', 36, True, GREEN, PP_ALIGN.CENTER)

tf6c = add_text_box(s6, 9.5, 2.8, 3.5, 3.5, '⚙️ PRO', 24, True, ORANGE)
add_bullet(tf6c, '', 6)
add_bullet(tf6c, 'Cotización', 18, WHITE, True)
add_bullet(tf6c, '4+ Workflows avanzados', 16, WHITE)
add_bullet(tf6c, 'Auditoría + monitoreo', 16, WHITE)

add_text_box(s6, 1, 6.3, 11, 0.8, 'No necesita invertir todo desde el inicio. Comience con lo que necesita hoy y crezca con Cresca.', 18, True, ACCENT, PP_ALIGN.CENTER)

# ===== SLIDE 7: Auditoría en 15 Minutos =====
s7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s7, WHITE)
add_text_box(s7, 1, 0.3, 11, 1, 'AUDITORÍA CRESCA EN 15 MINUTOS', 36, True, DARK, PP_ALIGN.CENTER)
add_text_box(s7, 1, 1.3, 11, 0.5, 'Diagnosticamos su negocio GRATIS para recomendarle el sistema correcto', 18, False, GRAY, PP_ALIGN.CENTER)

areas = [
    ('1️⃣  LEADS', '¿De dónde vienen\nsus clientes?', ACCENT),
    ('2️⃣  RESPUESTA', '¿Cuánto tardan en\nresponder?', GREEN),
    ('3️⃣  VENTAS', '¿Cómo cierran\nventas?', ACCENT),
    ('4️⃣  OPERACIONES', '¿Qué procesos\nconsumen tiempo?', GREEN),
    ('5️⃣  RETENCIÓN', '¿Cómo mantienen\ncontacto?', ACCENT),
]
for i, (title, desc, color) in enumerate(areas):
    x = 0.5 + i * 2.5
    tf = add_text_box(s7, x, 2.3, 2.2, 3, title, 18, True, color)
    add_bullet(tf, '', 6)
    add_bullet(tf, desc, 14, DARK)

add_text_box(s7, 1, 5.8, 11, 1, 'Resultado: Sabemos exactamente qué nivel necesita → STARTER, GROWTH o PRO AUTOMATION', 18, True, DARK, PP_ALIGN.CENTER)

# ===== SLIDE 8: Módulo de Reservas Universal =====
s8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s8, LIGHT_BG)
add_text_box(s8, 1, 0.5, 11, 1, 'MÓDULO DE RESERVAS — PARA CUALQUIER NEGOCIO', 32, True, DARK, PP_ALIGN.CENTER)
add_text_box(s8, 1, 1.5, 11, 0.8, 'Las reservas son una capacidad universal, disponible en Growth y Pro Automation', 18, False, GRAY, PP_ALIGN.CENTER)

tf8a = add_text_box(s8, 1, 2.8, 5, 3.5, '✅  Lo que incluye:', 20, True, DARK)
add_bullet(tf8a, '', 6)
add_bullet(tf8a, 'Sistema de reservas online', 16, DARK)
add_bullet(tf8a, 'Calendario de disponibilidad', 16, DARK)
add_bullet(tf8a, 'Notificaciones automáticas', 16, DARK)
add_bullet(tf8a, 'Sin comisiones de terceros', 16, DARK)
add_bullet(tf8a, 'Integrado con el sitio web', 16, DARK)

tf8b = add_text_box(s8, 7, 2.8, 5, 3.5, '🏢  Para quién:', 20, True, DARK)
add_bullet(tf8b, '', 6)
add_bullet(tf8b, 'Restaurantes y cafés', 16, DARK)
add_bullet(tf8b, 'Barberías y salones', 16, DARK)
add_bullet(tf8b, 'Hospedaje y Airbnb', 16, DARK)
add_bullet(tf8b, 'Clínicas y consultorios', 16, DARK)
add_bullet(tf8b, 'Cualquier negocio con citas', 16, DARK)

# ===== SLIDE 9: How We Work =====
s9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s9, DARK)
add_text_box(s9, 1, 0.5, 11, 1, 'CÓMO TRABAJAMOS', 36, True, WHITE, PP_ALIGN.CENTER)

tf9a = add_text_box(s9, 0.5, 2, 3.8, 4, '1️⃣  AUDITAMOS', 26, True, ACCENT)
add_bullet(tf9a, '', 8)
add_bullet(tf9a, '15 minutos de diagnóstico', 16, WHITE)
add_bullet(tf9a, 'Identificamos oportunidades', 16, WHITE)
add_bullet(tf9a, 'Recomendamos el nivel ideal', 16, WHITE)

tf9b = add_text_box(s9, 4.8, 2, 3.8, 4, '2️⃣  CONSTRUIMOS', 26, True, GREEN)
add_bullet(tf9b, '', 8)
add_bullet(tf9b, 'Implementamos su sistema', 16, WHITE)
add_bullet(tf9b, 'Configuramos workflows', 16, WHITE)
add_bullet(tf9b, 'Capacitamos a su equipo', 16, WHITE)

tf9c = add_text_box(s9, 9, 2, 3.8, 4, '3️⃣  OPTIMIZAMOS', 26, True, ACCENT)
add_bullet(tf9c, '', 8)
add_bullet(tf9c, 'Monitoreo continuo', 16, WHITE)
add_bullet(tf9c, 'Soporte y mantenimiento', 16, WHITE)
add_bullet(tf9c, 'Mejoras progresivas', 16, WHITE)

add_text_box(s9, 1, 6.2, 11, 0.8, 'No te dejamos solo después de la entrega. Somos tu socio tecnológico a largo plazo.', 18, True, ACCENT, PP_ALIGN.CENTER)

# ===== SLIDE 10: Payment Options =====
s10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s10, WHITE)
add_text_box(s10, 1, 0.5, 11, 1, 'OPCIONES DE PAGO FLEXIBLES', 36, True, DARK, PP_ALIGN.CENTER)

tf10 = add_text_box(s10, 1.5, 2, 10, 5, '', 18)
add_bullet(tf10, '📅  Mes a mes — sin compromiso, precio estándar', 20, DARK)
add_bullet(tf10, '', 8)
add_bullet(tf10, '💰  6 meses por adelantado — 1 mes gratis o descuento en setup', 20, DARK)
add_bullet(tf10, '', 8)
add_bullet(tf10, '🏆  12 meses por adelantado — 2 meses gratis o descuento en setup', 20, DARK)
add_bullet(tf10, '', 8)
add_bullet(tf10, '🤝  Plan Flex — Setup dividido: 50% hoy + 50% en 30 días', 20, DARK)
add_bullet(tf10, '', 8)
add_bullet(tf10, 'Empiece desde $9.99/mes — sin excusas para no comenzar', 20, GREEN, True)

# ===== SLIDE 11: Case Studies =====
s11 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s11, LIGHT_BG)
add_text_box(s11, 1, 0.5, 11, 1, 'RESULTADOS REALES', 36, True, DARK, PP_ALIGN.CENTER)
add_text_box(s11, 1, 1.5, 11, 0.5, 'Lo que Cresca hace por negocios como el tuyo', 18, False, GRAY, PP_ALIGN.CENTER)

tf11a = add_text_box(s11, 0.8, 2.5, 3.6, 4, '🏪 Emprendedor → STARTER', 20, True, DARK)
add_bullet(tf11a, '', 6)
add_bullet(tf11a, 'Antes: Solo WhatsApp, sin sistema', 14, GRAY)
add_bullet(tf11a, 'Después: Workflow Lite + plantillas', 14, DARK)
add_bullet(tf11a, 'Resultado: Respuesta automática 24/7 por $9.99/mes', 14, GREEN, True)

tf11b = add_text_box(s11, 4.8, 2.5, 3.6, 4, '🏨 Restaurante → GROWTH', 20, True, DARK)
add_bullet(tf11b, '', 6)
add_bullet(tf11b, 'Antes: Solo redes sociales, sin web', 14, GRAY)
add_bullet(tf11b, 'Después: Web + reservas + workflows', 14, DARK)
add_bullet(tf11b, 'Resultado: Reservas directas, más clientes', 14, GREEN, True)

tf11c = add_text_box(s11, 8.8, 2.5, 3.6, 4, '🏢 Empresa → PRO', 20, True, DARK)
add_bullet(tf11c, '', 6)
add_bullet(tf11c, 'Antes: Todo manual, Excel, papel', 14, GRAY)
add_bullet(tf11c, 'Después: 4+ workflows automáticos', 14, DARK)
add_bullet(tf11c, 'Resultado: 40% menos tiempo en admin', 14, GREEN, True)

add_text_box(s11, 1, 6.5, 11, 0.5, 'Casos basados en resultados típicos. Resultados específicos varían según el negocio.', 12, True, GRAY, PP_ALIGN.CENTER)

# ===== SLIDE 12: Next Steps =====
s12 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s12, DARK)
add_text_box(s12, 1, 0.8, 11, 1, '¿LISTO PARA CRECER?', 44, True, WHITE, PP_ALIGN.CENTER)

tf12 = add_text_box(s12, 2, 2.3, 9, 3.5, '', 20)
add_bullet(tf12, '1️⃣   Auditoría gratuita de 15 minutos', 22, WHITE)
add_bullet(tf12, '', 8)
add_bullet(tf12, '2️⃣   Te recomendamos el nivel ideal', 22, WHITE)
add_bullet(tf12, '', 8)
add_bullet(tf12, '3️⃣   Configuramos tu sistema', 22, WHITE)
add_bullet(tf12, '', 8)
add_bullet(tf12, '4️⃣   Tu negocio crece con Cresca', 22, WHITE)

add_text_box(s12, 2, 5.5, 9, 0.5, 'Desde $9.99/mes  ·  Sin contratos largos  ·  Empiece hoy', 20, True, GREEN, PP_ALIGN.CENTER)
add_text_box(s12, 2, 6.2, 9, 0.5, 'CRESCA  ·  Sistemas de Crecimiento Empresarial', 18, True, ACCENT, PP_ALIGN.CENTER)
add_text_box(s12, 2, 6.8, 9, 0.5, 'El Salvador, C.A.', 14, False, GRAY, PP_ALIGN.CENTER)

# Save
OUT = r'c:\Users\12132\Desktop\DigiFacil\DigiFacil\CRESCA_Pitch_Deck_UPDATED.pptx'
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
print("Done!")
